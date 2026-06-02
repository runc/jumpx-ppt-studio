"""Deck 导出（Phase 4b-1）：PDF + 逐页 PNG。

唯一渲染原语 = Playwright + Chromium（单机装一次）。我们的产物已是自包含
`index.html`（横向 flex deck，每页 .slide 为 100vw×100vh），所以：
- PDF：注入打印 CSS（deck 改 block、每页 1280×720 + break-after），page.pdf() 矢量输出。
- PNG：视口 1280×720@2x，逐页 transform 定位后整屏截图，打包 zip。

保真优先：直接渲染真实 HTML，CSS 主题 100% 保留。不依赖 LibreOffice。
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from playwright.async_api import async_playwright

import runs as RUN

# 16:9 导出基准（与 deck 的 100vw×100vh 对齐）
W, H = 1280, 720

# 打印态：把横向 flex deck 摊平成「每页一张、1280×720、分页」
_PRINT_CSS = f"""
@page {{ size: {W}px {H}px; margin: 0; }}
html, body {{ margin: 0 !important; padding: 0 !important; background: #fff !important;
  width: auto !important; height: auto !important; overflow: visible !important; }}
.deck {{ position: static !important; display: block !important;
  width: auto !important; height: auto !important; transform: none !important; }}
.slide {{ width: {W}px !important; height: {H}px !important;
  break-after: page; page-break-after: always; overflow: hidden; }}
.slide:last-child {{ break-after: auto; page-break-after: auto; }}
.slide-controls, .slide-index, [data-action], .notes {{ display: none !important; }}
"""

# 截图态：去过渡/去控件，逐页干净截屏
_SHOT_CSS = ".deck{transition:none !important}.slide-controls,.slide-index,[data-action]{display:none !important}"


def _export_dir(rid: str) -> Path | None:
    """run 的 export/ 目录（复用 runs 的安全校验）。"""
    html = RUN.index_html_path(rid)
    if html is None:
        return None
    d = html.parent / "export"
    d.mkdir(parents=True, exist_ok=True)
    return d


async def export_pdf(rid: str) -> Path | None:
    """渲染整本 deck 为矢量 PDF（每页一张幻灯片）。返回文件路径。"""
    html = RUN.index_html_path(rid)
    out_dir = _export_dir(rid)
    if html is None or out_dir is None:
        return None
    out = out_dir / f"{rid}.pdf"
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": W, "height": H})
        await page.goto(html.resolve().as_uri(), wait_until="networkidle")
        await page.add_style_tag(content=_PRINT_CSS)
        await page.emulate_media(media="print")
        await page.pdf(path=str(out), width=f"{W}px", height=f"{H}px",
                       print_background=True,
                       margin={"top": "0", "bottom": "0", "left": "0", "right": "0"})
        await browser.close()
    return out


async def _render_slide_pngs(html: Path) -> list[bytes]:
    """逐页截图（2x 清晰）→ 返回每页 PNG 字节。PNG/PPTX 共用。"""
    pngs: list[bytes] = []
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": W, "height": H},
                                      device_scale_factor=2)
        await page.goto(html.resolve().as_uri(), wait_until="networkidle")
        await page.add_style_tag(content=_SHOT_CSS)
        n = await page.eval_on_selector_all(".slide", "els => els.length")
        for i in range(max(1, n)):
            await page.evaluate(
                "(i) => { const d = document.getElementById('deck');"
                " if (d) d.style.transform = 'translateX(' + (-i * 100) + 'vw)'; }", i)
            await page.wait_for_timeout(140)
            pngs.append(await page.screenshot(type="png"))
        await browser.close()
    return pngs


async def export_png_zip(rid: str) -> Path | None:
    """逐页截图打包成 zip。返回 zip 路径。"""
    html = RUN.index_html_path(rid)
    out_dir = _export_dir(rid)
    if html is None or out_dir is None:
        return None
    pngs = await _render_slide_pngs(html)
    out = out_dir / f"{rid}-png.zip"
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, data in enumerate(pngs):
            zf.writestr(f"slide-{i + 1:02d}.png", data)
    return out


async def export_pptx(rid: str) -> Path | None:
    """图片忠实版 PPTX（Phase 4b-2）：每页整屏截图铺满一张 16:9 幻灯片。

    保真 100%（像素级还原 HTML 主题），不可编辑文字，但可在 PowerPoint/WPS
    移动/批注/放映。复用 PNG 截图链路。
    """
    import io

    from pptx import Presentation
    from pptx.util import Inches

    html = RUN.index_html_path(rid)
    out_dir = _export_dir(rid)
    if html is None or out_dir is None:
        return None
    pngs = await _render_slide_pngs(html)

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for data in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(io.BytesIO(data), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
    out = out_dir / f"{rid}.pptx"
    prs.save(str(out))
    return out
